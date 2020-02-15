from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

n = 0


def write_image(shape):
    global n
    image = shape.image
    # ---get image "file" contents---
    image_bytes = image.blob
    # ---make up a name for the file, e.g. 'image.jpg'---
    image_filename = 'image{:03d}.{}'.format(n, image.ext)
    n += 1
    print(image_filename)
    with open(image_filename, 'wb') as f:
        f.write(image_bytes)


def visitor(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            visitor(s)
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        write_image(shape)


def iter_picture_shapes(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            visitor(shape)


if __name__ == '__main__':
    iter_picture_shapes(Presentation("C:\Data\python\speech-test\samplepptx.pptx"))
